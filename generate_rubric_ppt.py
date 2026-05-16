from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()

# Standard blank slide layout
blank_slide_layout = prs.slide_layouts[6]
title_bullet_slide_layout = prs.slide_layouts[1]

# ========== SLIDE 1: TITLE SLIDE ==========
slide1 = prs.slides.add_slide(blank_slide_layout)

# College Name Header
tx_box1 = slide1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
tf1 = tx_box1.text_frame
tf1.word_wrap = True
p1 = tf1.add_paragraph()
p1.text = "NALLA MALLA REDDY ENGINEERING COLLEGE\nAutonomous Institution"
p1.font.bold = True
p1.font.size = Pt(28)
p1.font.color.rgb = RGBColor(128, 0, 0)
p1.alignment = PP_ALIGN.CENTER

# Project Title
tx_box_title = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
tf_title = tx_box_title.text_frame
tf_title.word_wrap = True
p_title = tf_title.add_paragraph()
p_title.text = "CodeSonar: AI-Powered Code Analysis & Visualization"
p_title.font.bold = True
p_title.font.size = Pt(40)
p_title.font.color.rgb = RGBColor(0, 51, 102)
p_title.alignment = PP_ALIGN.CENTER

p_subtitle = tf_title.add_paragraph()
p_subtitle.text = "Department of CSE (Artificial Intelligence & Machine Learning)\nMajor Project - Final Evaluation"
p_subtitle.font.size = Pt(20)
p_subtitle.alignment = PP_ALIGN.CENTER

# ========== SLIDE 2: INTRODUCTION ==========
slide2 = prs.slides.add_slide(title_bullet_slide_layout)
slide2.shapes.title.text = "1. Introduction"
tf2 = slide2.placeholders[1].text_frame
tf2.text = "Motivation: Rapidly shrinking developer onboarding time and technical debt visualization via AI."

p = tf2.add_paragraph()
p.text = "Project Context: CodeSonar is an enterprise-grade AI-powered codebase analysis platform."
p.level = 1

p = tf2.add_paragraph()
p.text = "Role: Acts as an automated \"Senior Engineer as a Service\" using large context window LLMs."
p.level = 1

p = tf2.add_paragraph()
p.text = "Capabilities:"
p.level = 0

p = tf2.add_paragraph()
p.text = "Identifies complete tech stack, bugs, and security vulnerabilities."
p.level = 1

p = tf2.add_paragraph()
p.text = "Scores code complexity and automatically drafts real-time system architectures."
p.level = 1

# ========== SLIDE 3: LITERATURE REVIEW ==========
slide3 = prs.slides.add_slide(title_bullet_slide_layout)
slide3.shapes.title.text = "2. Literature Review"
tf3 = slide3.placeholders[1].text_frame
tf3.text = "Current State of Code Analysis:"

p = tf3.add_paragraph()
p.text = "Traditional Static Analyzers (SonarQube, ESLint): Effective at syntax rule enforcing, but fundamentally lack semantic or architectural-level understanding."
p.level = 1

p = tf3.add_paragraph()
p.text = "Generative AI Assistants (GitHub Copilot, ChatGPT): Powerful at localized function generation, but lack overarching repository structural map comprehension."
p.level = 1

p = tf3.add_paragraph()
p.text = "Recent Advancements: Emergence of massive context-window LLMs (e.g., Gemini 2.5 Flash) allows for 'Context Stuffing' without the need for convoluted Vector Databases (RAG approaches)."
p.level = 1

# ========== SLIDE 4: RESEARCH GAP ==========
slide4 = prs.slides.add_slide(title_bullet_slide_layout)
slide4.shapes.title.text = "3. Research Gap"
tf4 = slide4.placeholders[1].text_frame
tf4.text = "Bridging the Divide: There is a prominent gap between isolated static linting tools and high-level conversational AI."

p = tf4.add_paragraph()
p.text = "Missing Visual Integration: Current tools do not organically translate codebase dependencies into immediate visual architectures (flowcharts/Mermaid.js)."
p.level = 1

p = tf4.add_paragraph()
p.text = "Lack of Unification: No easily accessible platform unifies bug detection, cyclomatic complexity scoring, and file-by-file text explanation into a single, cohesive presentation layer."
p.level = 1

# ========== SLIDE 5: PROBLEM STATEMENT ==========
slide5 = prs.slides.add_slide(title_bullet_slide_layout)
slide5.shapes.title.text = "4. Problem Statement"
tf5 = slide5.placeholders[1].text_frame
tf5.text = "The Overwhelming Complexity of Large Repositories"

p = tf5.add_paragraph()
p.text = "As software ecosystems scale, manual codebase onboarding becomes incredibly resource-intensive and prone to friction."
p.level = 1

p = tf5.add_paragraph()
p.text = "Legacy code bases often suffer from \"Knowledge Silos\" where documentation is either missing, disjointed, or hopelessly outdated."
p.level = 1

p = tf5.add_paragraph()
p.text = "Precise Challenge: We need an automated, dynamic system that ingests entire repositories straight from source and outputs interactive, human-readable insights mapping dependencies and security risks automatically."
p.level = 1

# ========== SLIDE 6: OBJECTIVES (INNOVATIVE & NOVEL) ==========
slide6 = prs.slides.add_slide(title_bullet_slide_layout)
slide6.shapes.title.text = "5. Objectives (Innovative & Novel)"
tf6 = slide6.placeholders[1].text_frame
tf6.text = "Highly Innovative Goals of CodeSonar:"

p = tf6.add_paragraph()
p.text = "Automated Source Ingestion: Recursively fetch and filter GitHub repos via Octokit API dynamically."
p.level = 1

p = tf6.add_paragraph()
p.text = "Comprehensive Contextual Assessment: Score complexity mapping and flag architectural anti-patterns using custom AI prompt-agents."
p.level = 1

p = tf6.add_paragraph()
p.text = "Dynamic Architectural Rendering: Translate AI logic derivations directly into executable Mermaid.js charts mapping internal data flow organically."
p.level = 1

p = tf6.add_paragraph()
p.text = "Context-Aware Agent Chatbot: Provision a real-time conversational layer capable of answering hyper-specific design choices using the stuffed repository context."
p.level = 1

# Save Presentation
output_path = 'r:/dup/codesonar/CodeSonar_Final_Evaluation.pptx'
prs.save(output_path)
print(f"Presentation saved successfully at {output_path}")
